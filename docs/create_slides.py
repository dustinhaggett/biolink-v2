"""Generate BioLink final presentation slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x0B, 0x14, 0x1F)
TEAL = RGBColor(0x00, 0x60, 0x6D)
LIGHT_TEAL = RGBColor(0x26, 0x79, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC0)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
RED = RGBColor(0xC4, 0x32, 0x0A)
AMBER = RGBColor(0xB4, 0x53, 0x09)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullet_slide(slide, left, top, width, height, items, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)

        if isinstance(item, tuple):
            # (bold_part, rest)
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(size)
            run1.font.color.rgb = ACCENT_BLUE
            run1.font.bold = True
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(size)
            run2.font.color.rgb = color
        else:
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = color
    return tf


def add_card(slide, left, top, width, height, title, body, title_color=ACCENT_BLUE):
    """Add a rounded card with title and body text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x14, 0x20, 0x30)
    shape.line.color.rgb = RGBColor(0x1E, 0x33, 0x4A)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = LIGHT_GRAY


# ═══════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_text(slide, 1, 1.5, 11, 1.2, "BioLink", size=52, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 2.7, 11, 0.8, "AI-Powered Drug Repurposing Discovery", size=28, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.5, "Turning a research model into a tool that patients, researchers, and clinicians can actually use", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.4, "AAI 595 — Applied Machine Learning  |  Dr. Tao Han", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6.0, 11, 0.4, "Dustin Haggett  |  Kera Prosper  |  Esume", size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "The Problem", size=36, color=ACCENT_BLUE, bold=True)
add_text(slide, 0.8, 1.2, 11, 0.5, "Drug development is broken. Can AI help fix it?", size=18, color=LIGHT_GRAY)

add_card(slide, 0.8, 2.0, 3.6, 1.8, "10-15 Years", "Average time to bring a new drug to market from discovery to FDA approval")
add_card(slide, 4.8, 2.0, 3.6, 1.8, "$2.6 Billion", "Average cost per approved drug, including failures (DiMasi et al., 2016)")
add_card(slide, 8.8, 2.0, 3.6, 1.8, "90% Failure Rate", "Nine out of ten drug candidates fail in clinical trials")

add_text(slide, 0.8, 4.2, 11.5, 0.5, "Drug repurposing offers a faster path:", size=20, color=WHITE, bold=True)
add_bullet_slide(slide, 0.8, 4.8, 11, 2.2, [
    "7,000+ FDA-approved drugs already have established safety profiles",
    "Repurposed drugs can skip years of preclinical testing",
    "But identifying candidates requires massive literature review and domain expertise",
    "No accessible tool exists for exploring drug-disease connections grounded in evidence",
], size=16)

# ═══════════════════════════════════════════════════
# SLIDE 3: Our Solution
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "Our Solution: BioLink", size=36, color=ACCENT_BLUE, bold=True)
add_text(slide, 0.8, 1.2, 11, 0.5, "A knowledge graph AI that makes drug repurposing accessible to everyone", size=18, color=LIGHT_GRAY)

add_card(slide, 0.8, 2.2, 5.5, 1.4, "Search by Disease", "\"What drugs might treat Lyme Disease?\" — Ranks 7,163 drugs by predicted relevance with calibrated confidence scores")
add_card(slide, 6.8, 2.2, 5.5, 1.4, "Search by Drug", "\"I'm taking Metformin — what else might it help?\" — Scores all 2,525 diseases against a single drug")
add_card(slide, 0.8, 4.0, 3.5, 1.5, "Evidence Grounding", "Every prediction is checked against published literature, clinical trials, and FDA data — not just model confidence")
add_card(slide, 4.7, 4.0, 3.5, 1.5, "Verdict System", "Evidence Supports / Conflicts / Standard-of-Care / Insufficient — so users see when evidence disagrees with the model")
add_card(slide, 8.6, 4.0, 3.5, 1.5, "Actionable Output", "Clinical trial links, PDF reports, comparison tables, follow-up Q&A — built for real-world use")

add_text(slide, 0.8, 5.9, 11, 0.5, "Always with disclaimers: BioLink is a hypothesis generator, not medical advice.", size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 4: How It Works (Architecture)
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "How It Works", size=36, color=ACCENT_BLUE, bold=True)

# Pipeline steps as cards
steps = [
    ("1. Input", "User types disease\nor drug name"),
    ("2. Entity Resolution", "Claude Haiku maps\nfree-text to CTD entity"),
    ("3. Model Scoring", "MLP scores 7,163 drugs\n(AUC = 0.947)"),
    ("4. Calibration", "Temperature scaling\nfor honest probabilities"),
    ("5. Enrichment", "PubMed + FDA + Trials\n(parallel async)"),
    ("6. Evidence", "Perplexity searches\npublished literature"),
]

for i, (title, body) in enumerate(steps):
    x = 0.5 + i * 2.05
    add_card(slide, x, 1.6, 1.9, 1.6, title, body)

# Model details
add_text(slide, 0.8, 3.6, 11, 0.5, "The AI Model", size=22, color=WHITE, bold=True)
add_card(slide, 0.8, 4.2, 3.6, 2.5, "MLP Classifier",
    "Input: 800-dim feature vector\n"
    "[drug, disease, |diff|, product]\n\n"
    "Hidden: 256 neurons\n"
    "BatchNorm + ReLU + Dropout\n\n"
    "Output: calibrated probability")
add_card(slide, 4.8, 4.2, 3.6, 2.5, "BioWordVec Embeddings",
    "200-dim vectors trained on\n"
    "PubMed + MIMIC-III\n\n"
    "Captures biomedical semantic\n"
    "relationships between terms\n\n"
    "7,163 drugs + 2,525 diseases\n"
    "pre-embedded at startup")
add_card(slide, 8.8, 4.2, 3.6, 2.5, "Training Data",
    "Comparative Toxicogenomics\n"
    "Database (CTD)\n\n"
    "~380,000 drug-disease pairs\n"
    "with known associations\n\n"
    "Test AUC: 0.947")

# ═══════════════════════════════════════════════════
# SLIDE 5: Key Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "Key Features", size=36, color=ACCENT_BLUE, bold=True)

features = [
    ("Evidence Verdicts", "Supports / Conflicts /\nStandard-of-Care /\nInsufficient", GREEN),
    ("Evidence Quality", "RCT / Human Study /\nPreclinical / Case Report /\nTheoretical", ACCENT_BLUE),
    ("Clinical Trials", "Live lookup from\nClinicalTrials.gov with\ndirect enrollment links", TEAL),
    ("Pathway Chains", "Drug -> Target ->\nPathway -> Effect\nvisualized per result", LIGHT_TEAL),
    ("Compare Mode", "Select multiple drugs\nfor side-by-side\ncomparison table", AMBER),
    ("Follow-up Q&A", "Ask questions about\nany result, grounded\nin published evidence", ACCENT_BLUE),
    ("Drug Interactions", "Warns when a candidate\ninteracts with standard\ntreatments", RED),
    ("Export", "PDF reports + CSV\nfor sharing with\nhealthcare providers", LIGHT_GRAY),
]

for i, (title, body, color) in enumerate(features):
    row = i // 4
    col = i % 4
    x = 0.5 + col * 3.1
    y = 1.5 + row * 2.6
    add_card(slide, x, y, 2.9, 2.2, title, body, title_color=color)

# ═══════════════════════════════════════════════════
# SLIDE 6: API Integration
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "Multi-API Integration", size=36, color=ACCENT_BLUE, bold=True)
add_text(slide, 0.8, 1.2, 11, 0.5, "5 APIs working together — 3 free, 2 paid (cheap)", size=18, color=LIGHT_GRAY)

apis = [
    ("Claude Haiku (Anthropic)", "Entity resolution + explanations", "API key", "~$0.001/query", ACCENT_BLUE),
    ("Perplexity Sonar", "Evidence search + follow-up Q&A", "API key", "~$0.005/query", ACCENT_BLUE),
    ("PubMed E-utilities", "Publication co-occurrence counts", "None", "Free", GREEN),
    ("OpenFDA", "Drug approval status lookup", "None", "Free", GREEN),
    ("ClinicalTrials.gov v2", "Active/completed trial finder", "None", "Free", GREEN),
]

# Table header
add_text(slide, 0.8, 2.0, 4, 0.4, "API", size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, 4.8, 2.0, 3.5, 0.4, "Purpose", size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, 8.3, 2.0, 1.5, 0.4, "Auth", size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, 9.8, 2.0, 2, 0.4, "Cost", size=14, color=ACCENT_BLUE, bold=True)

for i, (name, purpose, auth, cost, color) in enumerate(apis):
    y = 2.5 + i * 0.55
    add_text(slide, 0.8, y, 4, 0.4, name, size=14, color=WHITE, bold=True)
    add_text(slide, 4.8, y, 3.5, 0.4, purpose, size=14, color=LIGHT_GRAY)
    add_text(slide, 8.3, y, 1.5, 0.4, auth, size=14, color=LIGHT_GRAY)
    add_text(slide, 9.8, y, 2, 0.4, cost, size=14, color=color)

add_text(slide, 0.8, 5.5, 11, 0.5, "All enrichment runs in parallel using Python asyncio — minimal added latency", size=16, color=LIGHT_GRAY)
add_text(slide, 0.8, 6.0, 11, 0.5, "App works without API keys (fuzzy matching fallback) — keys unlock full LLM features", size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 7: Innovation — Evidence vs Model
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "Innovation: When AI Disagrees with Evidence", size=32, color=ACCENT_BLUE, bold=True)
add_text(slide, 0.8, 1.2, 11, 0.5, "The most useful feature isn't the model — it's showing when the model is wrong", size=18, color=LIGHT_GRAY)

add_card(slide, 0.8, 2.2, 5.5, 2.2, "Cyclosporine for Lyme Disease",
    "Model says: 99% Strong confidence\n"
    "Evidence says: CONFLICTS\n\n"
    "Cyclosporine is an immunosuppressant.\n"
    "Lyme Disease is a bacterial infection.\n"
    "Immunosuppression would worsen it.\n\n"
    "The model detects a real biological\n"
    "relationship but can't distinguish\n"
    "\"interacts with\" from \"would help treat.\"", title_color=RED)

add_card(slide, 6.8, 2.2, 5.5, 2.2, "Doxycycline for Lyme Disease",
    "Model says: 99% Strong confidence\n"
    "Evidence says: STANDARD-OF-CARE\n\n"
    "Doxycycline is already the first-line\n"
    "antibiotic treatment for early Lyme.\n"
    "Multiple RCTs, CDC guidelines.\n\n"
    "Model and evidence align perfectly.\n"
    "This is the actionable result.", title_color=GREEN)

add_text(slide, 0.8, 4.8, 11.5, 0.8,
    "Key insight: The gap between model confidence and real-world evidence IS the information.",
    size=20, color=WHITE, bold=True)
add_text(slide, 0.8, 5.5, 11.5, 0.8,
    "No other drug repurposing tool shows both signals. Users see exactly where to trust the model and where to be skeptical.",
    size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 8: Challenges & Solutions
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "Challenges & Solutions", size=36, color=ACCENT_BLUE, bold=True)

challenges = [
    ("Model vs. Evidence Mismatch",
     "Model gives high confidence to drugs that\nevidence shows are harmful",
     "Integrated Perplexity evidence search with\nverdict badges showing agreement/conflict"),
    ("API Cost at Scale",
     "Perplexity costs ~$0.01 per call;\n20 results = $0.20 per query",
     "Tiered approach: evidence for top 5,\nfree APIs (PubMed/FDA/Trials) for all 20"),
    ("Free-Text Entity Resolution",
     "Users type 'heart attack' but model\nneeds 'Myocardial Infarction'",
     "Claude Haiku + difflib fuzzy fallback;\nclarification prompt for low confidence"),
    ("Reverse Search Architecture",
     "Model only scored drugs for diseases,\nnot diseases for drugs",
     "Pre-cache 2,525 disease embeddings;\nmirror score_all_drugs() method"),
]

for i, (title, problem, solution) in enumerate(challenges):
    y = 1.5 + i * 1.4
    add_text(slide, 0.8, y, 3.5, 0.4, title, size=16, color=ACCENT_BLUE, bold=True)
    add_text(slide, 4.5, y, 4, 0.6, problem, size=13, color=LIGHT_GRAY)
    add_text(slide, 8.8, y, 4, 0.6, solution, size=13, color=WHITE)

# Column headers
add_text(slide, 4.5, 1.1, 4, 0.4, "Challenge", size=14, color=RED, bold=True)
add_text(slide, 8.8, 1.1, 4, 0.4, "Solution", size=14, color=GREEN, bold=True)

# ═══════════════════════════════════════════════════
# SLIDE 9: AI Tools Usage
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "AI Tools Usage", size=36, color=ACCENT_BLUE, bold=True)
add_text(slide, 0.8, 1.2, 11, 0.5, "AI was both the product AND the development tool", size=18, color=LIGHT_GRAY)

add_card(slide, 0.8, 2.0, 5.5, 2.0, "Development Tools",
    "Claude Code — Architecture design, code generation,\n"
    "debugging, all 9 features, documentation\n\n"
    "GitHub Copilot — Code completion\n\n"
    "Result: Built a 10-file, 2000+ line application\n"
    "with 5 API integrations in weeks, not months")

add_card(slide, 6.8, 2.0, 5.5, 2.0, "Runtime AI Components",
    "Claude Haiku — Entity resolution (natural language\n"
    "to standardized medical terms)\n\n"
    "Perplexity Sonar — Live evidence search, structured\n"
    "verdict extraction, follow-up Q&A\n\n"
    "Result: Every prediction grounded in real evidence")

add_text(slide, 0.8, 4.5, 11, 0.8,
    "AI tools didn't just help write code — they identified architectural patterns, caught bugs the model "
    "introduced, and suggested the evidence verdict system that became the app's most distinctive feature.",
    size=16, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 10: Impact & Future
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.8, 0.5, 11, 0.7, "Impact & Future Work", size=36, color=ACCENT_BLUE, bold=True)

add_card(slide, 0.8, 1.5, 3.5, 2.5, "For Patients",
    "Explore drug candidates for\nconditions with limited options\n\n"
    "Bring evidence-backed PDF\nreports to their doctors\n\n"
    "Ask follow-up questions\nabout specific results")
add_card(slide, 4.7, 1.5, 3.5, 2.5, "For Researchers",
    "Screen candidates across\nthousands of drug-disease pairs\n\n"
    "Direct links to active\nclinical trials\n\n"
    "Batch mode for\nbulk hypothesis generation")
add_card(slide, 8.6, 1.5, 3.5, 2.5, "For Clinicians",
    "Compare candidates\nside-by-side\n\n"
    "Drug interaction warnings\nfor standard-of-care\n\n"
    "Evidence quality ratings\n(RCT vs case report)")

add_text(slide, 0.8, 4.4, 11, 0.5, "Future Work", size=22, color=WHITE, bold=True)
add_bullet_slide(slide, 0.8, 5.0, 11, 2, [
    "Deploy to Hugging Face Spaces for public access",
    "Retrain model with updated CTD data and graph neural network architecture",
    "Add drug-drug interaction database (DrugBank) for deeper safety checks",
    "Integrate patient-specific factors (age, comorbidities) for personalized ranking",
], size=16)

# ═══════════════════════════════════════════════════
# SLIDE 11: Live Demo
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 2.5, 11, 1.2, "Live Demo", size=52, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.8, "biolink-v2  |  streamlit run app.py", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 12: Thank You
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 2.0, 11, 1.0, "Thank You", size=48, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.2, 11, 0.6, "Questions?", size=28, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11, 0.4, "github.com/dustinhaggett/biolink-v2", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.2, 11, 0.4, "BioLink is a hypothesis-generation tool for research purposes only. Not medical advice.", size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save
output = "/Users/berlin/biolink_v2/docs/BioLink_Final_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
